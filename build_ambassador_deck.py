#!/usr/bin/env python3
"""
PISES New Campus – Ambassador Highlights Deck Generator
Produces a 3-slide executive briefing PowerPoint for Embassy / SMC review.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ─────────────────────────────────────────────────────────────────────────────
# COLOUR PALETTE (Pakistan flag inspired + institutional)
# ─────────────────────────────────────────────────────────────────────────────
DARK_GREEN   = RGBColor(0x01, 0x41, 0x1C)   # Pakistan green
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GOLD         = RGBColor(0xC4, 0x9A, 0x2A)   # accent gold
LIGHT_GREEN  = RGBColor(0x4C, 0xAF, 0x50)
LIGHT_GREY   = RGBColor(0xF5, 0xF5, 0xF5)
MID_GREY     = RGBColor(0xBD, 0xBD, 0xBD)
DARK_GREY    = RGBColor(0x42, 0x42, 0x42)
BLACK        = RGBColor(0x00, 0x00, 0x00)
HEADER_BG    = RGBColor(0x01, 0x41, 0x1C)
ROW_ALT      = RGBColor(0xE8, 0xF5, 0xE9)
ACCENT_RED   = RGBColor(0xC6, 0x28, 0x28)
BLUE_ACCENT  = RGBColor(0x1B, 0x5E, 0x20)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

# ─────────────────────────────────────────────────────────────────────────────
# HELPER FUNCTIONS
# ─────────────────────────────────────────────────────────────────────────────
def add_bg(slide, color=DARK_GREEN):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.0):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=11,
                      color=BLACK, font_name='Calibri', alignment=PP_ALIGN.LEFT,
                      bold_first=False, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.font.bold = (bold_first and i == 0)
        p.space_after = Pt(2)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table

def style_header_cell(cell, text, font_size=9):
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = HEADER_BG
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def style_data_cell(cell, text, font_size=8, bold=False, alignment=PP_ALIGN.CENTER,
                    fill_color=None, font_color=BLACK):
    cell.text = str(text)
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color
        p.font.name = 'Calibri'
        p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_kpi_card(slide, left, top, width, height, label, value, sub="",
                 bg_color=WHITE, value_color=DARK_GREEN, label_color=DARK_GREY):
    shape = add_rect(slide, left, top, width, height, bg_color)
    shape.shadow.inherit = False
    # Value
    add_text_box(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.55),
                 value, font_size=22, bold=True, color=value_color, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.15), top + Inches(0.62), width - Inches(0.3), Inches(0.35),
                 label, font_size=9, bold=False, color=label_color, alignment=PP_ALIGN.CENTER)
    if sub:
        add_text_box(slide, left + Inches(0.15), top + Inches(0.88), width - Inches(0.3), Inches(0.25),
                     sub, font_size=7, bold=False, color=MID_GREY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – DESIGN FRAMEWORK & REGULATORY BASIS
# ═══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, WHITE)

# Top banner
add_rect(slide1, Inches(0), Inches(0), Inches(16), Inches(1.15), DARK_GREEN)
add_text_box(slide1, Inches(0.5), Inches(0.15), Inches(10), Inches(0.55),
             "PISES NEW CAMPUS  |  AMBASSADOR HIGHLIGHTS DECK",
             font_size=22, bold=True, color=WHITE, font_name='Calibri')
add_text_box(slide1, Inches(0.5), Inches(0.65), Inches(10), Inches(0.4),
             "Pakistan International School (English Section), Riyadh  |  Al Safa Plot  |  25,000 m\u00b2",
             font_size=11, bold=False, color=GOLD, font_name='Calibri')
# Slide tag
add_text_box(slide1, Inches(12.5), Inches(0.25), Inches(3), Inches(0.5),
             "SLIDE 1 OF 3", font_size=10, bold=True, color=GOLD, alignment=PP_ALIGN.RIGHT)
add_text_box(slide1, Inches(12.5), Inches(0.55), Inches(3), Inches(0.4),
             "Design Framework & Regulatory Basis", font_size=10, bold=False, color=WHITE, alignment=PP_ALIGN.RIGHT)

# ── SECTION A: KPI CARDS ──
kpi_y = Inches(1.4)
kpi_h = Inches(1.15)
kpi_w = Inches(2.15)
gap = Inches(0.22)
start_x = Inches(0.5)

kpis = [
    ("CATEGORY A", "TBC Classification", "Major Cities (Riyadh)"),
    ("4.8 m\u00b2/student", "Education Complex Baseline", "(K+Elem+Inter+Sec)/4"),
    ("25 students/class", "Max Classroom Capacity", "TBC Mandated Limit"),
    ("~304 classrooms", "7,000-Student Model", "Gender-separated G2-G12"),
    ("~52,400 m\u00b2", "Total Built-Up Area (BUA)", "NET 33,983 m\u00b2 \u00d7 Grossing"),
    ("B + G + 2", "Building Configuration", "1 Basement + 3 Above Grade"),
]

for i, (val, label, sub) in enumerate(kpis):
    x = start_x + i * (kpi_w + gap)
    add_kpi_card(slide1, x, kpi_y, kpi_w, kpi_h, label, val, sub,
                 bg_color=RGBColor(0xF1, 0xF8, 0xE9))

# ── SECTION B: TBC vs Non-TBC Requirements Matrix ──
sec_b_y = Inches(2.85)
add_text_box(slide1, Inches(0.5), sec_b_y, Inches(8), Inches(0.35),
             "TBC-MANDATED vs NON-TBC DESIGN STANDARDS  |  CATEGORY A  (Riyadh)",
             font_size=12, bold=True, color=DARK_GREEN)

# Left table: TBC Area Per Student Standards
tbl1_top = sec_b_y + Inches(0.4)
tbl1 = add_table(slide1, 10, 5, Inches(0.5), tbl1_top, Inches(7.2), Inches(3.6))

# Headers
headers1 = ["Education Level", "Area/Student (m\u00b2)", "Example Net (m\u00b2)", "Facility", "Ratio Rule"]
for j, h in enumerate(headers1):
    style_header_cell(tbl1.cell(0, j), h, font_size=8)

# Data rows for TBC Category A
tbc_data = [
    ["Nursery (1m-1yr)", "1.4 (activity) + 0.7 (bed)", "45.0 per room", "Activity Room + Bedroom", "Per TBC manual"],
    ["Nursery (1yr-3yr)", "1.8 (activity) + 0.9 (bed)", "45.0 + 22.5", "Activity + Bedroom + Feeding", "Per TBC manual"],
    ["Kindergarten", "2.5", "62.5", "Classroom + Toilet + Court", "Court = \u00bd class area"],
    ["Elementary", "1.3 (class) / 1.9 (lab)", "43.12 / 60.1", "Class + Lab + LRC + Art + MPR", "1 Lab per 10 classes"],
    ["Intermediate", "1.4 (class) / 2.0 (lab)", "45.90 / 62.96", "Class + Labs + LRC + Art + MPR", "1 Lab per 10 classes"],
    ["Secondary", "1.5 (class) / 2.2 (lab)", "48.12 / 69.9", "Class + Labs + LRC + Art + MPR", "1 Lab per 10 classes"],
    ["Science Lab (Elem)", "1.9", "60.1", "General Science Lab", "1 per 10 classrooms"],
    ["Computer Lab (Elem)", "1.9", "60.1", "Computer & Languages Lab", "TBC unit / qty planned"],
    ["Learning Resource Ctr", "2.5", "75.1", "Library / LRC", "Librarian 8 m\u00b2 min"],
]
for i, row in enumerate(tbc_data):
    bg = ROW_ALT if i % 2 == 0 else WHITE
    for j, val in enumerate(row):
        style_data_cell(tbl1.cell(i+1, j), val, font_size=7.5, fill_color=bg,
                       alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

# Set column widths
col_widths_1 = [Inches(1.6), Inches(1.5), Inches(1.2), Inches(1.6), Inches(1.3)]
for j, w in enumerate(col_widths_1):
    tbl1.columns[j].width = w

# Right panel: Regulatory Hierarchy & Grossing Factors
panel_x = Inches(8.1)
add_text_box(slide1, panel_x, sec_b_y, Inches(7.5), Inches(0.35),
             "REGULATORY HIERARCHY & NET\u2192GROSS METHODOLOGY",
             font_size=12, bold=True, color=DARK_GREEN)

# Regulatory hierarchy box
rh_top = sec_b_y + Inches(0.45)
add_rect(slide1, panel_x, rh_top, Inches(3.6), Inches(2.0), RGBColor(0xF1, 0xF8, 0xE9))
reg_lines = [
    "REGULATORY HIERARCHY (Order of Authority):",
    "1. Saudi Building Code (SBC) \u2014 Mandatory",
    "2. Civil Defense Regulations \u2014 Mandatory",
    "3. Municipal / Balady Requirements \u2014 Mandatory",
    "4. TBC Category A Guidelines \u2014 Programmatic",
    "5. International Codes (IBC/NFPA) \u2014 Reference",
    "",
    "\u25b6 Where conflict: more stringent governs",
    "\u25b6 Occupancy: Group E (Educational)",
    "\u25b6 Risk Category III (>250 occupants)",
]
add_multiline_box(slide1, panel_x + Inches(0.12), rh_top + Inches(0.08),
                  Inches(3.4), Inches(1.85), reg_lines, font_size=8,
                  color=DARK_GREY, bold_first=True, line_spacing=1.3)

# Grossing factors box
gf_x = panel_x + Inches(3.85)
add_rect(slide1, gf_x, rh_top, Inches(3.55), Inches(2.0), RGBColor(0xFFF8, 0xE1, 0x00) if False else RGBColor(0xFF, 0xF8, 0xE1))
gf_lines = [
    "NET \u2192 GROSS MULTIPLIERS:",
    "Academic / Admin / SEN:    NET \u00d7 1.45",
    "High-Service (Labs/Sports): NET \u00d7 1.65",
    "Ops / Back-of-House:       NET \u00d7 1.55",
    "",
    "EDUCATION COMPLEX FORMULA:",
    "(K + Elem + Inter + Sec) \u00f7 4 = 4.8 m\u00b2/student",
    "(5.0 + 4.4 + 4.7 + 5.25) \u00f7 4 = 4.8375",
    "",
    "Land Baseline: 4.8 m\u00b2 \u00d7 7,000 = 33,600 m\u00b2",
]
add_multiline_box(slide1, gf_x + Inches(0.12), rh_top + Inches(0.08),
                  Inches(3.35), Inches(1.85), gf_lines, font_size=8,
                  color=DARK_GREY, bold_first=True, line_spacing=1.3)

# Non-TBC items mini-table
ntbc_top = rh_top + Inches(2.15)
add_rect(slide1, panel_x, ntbc_top, Inches(7.4), Inches(0.3), DARK_GREEN)
add_text_box(slide1, panel_x + Inches(0.1), ntbc_top + Inches(0.03), Inches(7.2), Inches(0.25),
             "NON-TBC ITEMS (Best Practice / Institutional Planning Assumptions)",
             font_size=8, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

ntbc_items_top = ntbc_top + Inches(0.35)
ntbc_tbl = add_table(slide1, 5, 4, panel_x, ntbc_items_top, Inches(7.4), Inches(1.5))
ntbc_headers = ["Category", "Facilities", "Net Area Drivers", "Basis"]
for j, h in enumerate(ntbc_headers):
    style_header_cell(ntbc_tbl.cell(0, j), h, font_size=7)

ntbc_data = [
    ["SEN & Wellbeing", "Resource Rooms, Therapy, Sensory, Counselling", "736 m\u00b2 NET", "International Best Practice"],
    ["Food Services", "2\u00d7 Dining Halls, Kitchen, Servery, Cold Store", "3,380 m\u00b2 NET", "Multi-shift dining model"],
    ["Sports & PE", "2\u00d7 Sports Halls, Pool (25m), Changing", "3,981 m\u00b2 NET", "Institutional standard"],
    ["Auditorium/Commons", "300-seat Auditorium, 2,000 m\u00b2 Atrium, Exam Hall", "4,070 m\u00b2 NET", "Campus life / events"],
]
for i, row in enumerate(ntbc_data):
    bg = ROW_ALT if i % 2 == 0 else WHITE
    for j, val in enumerate(row):
        style_data_cell(ntbc_tbl.cell(i+1, j), val, font_size=7, fill_color=bg,
                       alignment=PP_ALIGN.LEFT)

ntbc_col_widths = [Inches(1.4), Inches(2.6), Inches(1.3), Inches(2.1)]
for j, w in enumerate(ntbc_col_widths):
    ntbc_tbl.columns[j].width = w

# Footer
add_rect(slide1, Inches(0), Inches(8.55), Inches(16), Inches(0.45), DARK_GREEN)
add_text_box(slide1, Inches(0.5), Inches(8.58), Inches(10), Inches(0.35),
             "CONFIDENTIAL  |  Pakistan International School (English Section), Riyadh  |  Basis of Design v0.4",
             font_size=8, bold=False, color=GOLD)
add_text_box(slide1, Inches(12), Inches(8.58), Inches(3.5), Inches(0.35),
             "Prepared for Ambassador / SMC Briefing",
             font_size=8, bold=False, color=WHITE, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – CAPACITY SCENARIOS & AREA COMPUTATION
# ═══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

# Top banner
add_rect(slide2, Inches(0), Inches(0), Inches(16), Inches(1.15), DARK_GREEN)
add_text_box(slide2, Inches(0.5), Inches(0.15), Inches(10), Inches(0.55),
             "CAPACITY SCENARIOS  |  5,500 / 6,000 / 7,000 STUDENTS",
             font_size=22, bold=True, color=WHITE)
add_text_box(slide2, Inches(0.5), Inches(0.65), Inches(10), Inches(0.4),
             "Proportional Scaling from Actual Enrollment (5,263 current) to Design Targets  |  NET-First Model",
             font_size=11, bold=False, color=GOLD)
add_text_box(slide2, Inches(12.5), Inches(0.25), Inches(3), Inches(0.5),
             "SLIDE 2 OF 3", font_size=10, bold=True, color=GOLD, alignment=PP_ALIGN.RIGHT)
add_text_box(slide2, Inches(12.5), Inches(0.55), Inches(3), Inches(0.4),
             "Scenario Comparison & Area Build-Up", font_size=10, bold=False, color=WHITE, alignment=PP_ALIGN.RIGHT)

# ── Current Enrollment Summary ──
enroll_y = Inches(1.35)
add_text_box(slide2, Inches(0.5), enroll_y, Inches(6), Inches(0.3),
             "CURRENT ENROLLMENT SNAPSHOT  (Session 2024-25, as of 31 March 2025)",
             font_size=11, bold=True, color=DARK_GREEN)

enroll_tbl_top = enroll_y + Inches(0.35)
enroll_tbl = add_table(slide2, 5, 8, Inches(0.5), enroll_tbl_top, Inches(7.2), Inches(1.65))

e_headers = ["Segment", "Grades", "Students", "Boys", "Girls", "Sections", "Avg/Section", "Classrooms\n@25 cap"]
for j, h in enumerate(e_headers):
    style_header_cell(enroll_tbl.cell(0, j), h, font_size=7)

# Actual enrollment data computed from XLS
# Nursery: 122, Reception: 296, KG: 383 = 801 Early Years
# G1: 394, G2: 417, G3: 398, G4: 430 = 1,639 Primary (incl G1)
# G5: 422, G6: 468, G7: 429, G8: 381, G9: 351 = 2,051 Intermediate
# G10: 330, G11: 253, G12: 189 = 772 Secondary
# Total: 5,263 (Boys: 2,690, Girls: 2,573)

enrollment_rows = [
    ["Early Years", "Nursery-KG", "801", "415", "386", "55", "14.6", "33"],
    ["Primary", "G1\u2013G4", "1,639", "839", "800", "82", "20.0", "66"],
    ["Intermediate", "G5\u2013G9", "2,051", "1,042", "1,009", "104", "19.7", "83"],
    ["Secondary", "G10\u2013G12", "772", "394", "378", "40", "19.3", "31"],
]
for i, row in enumerate(enrollment_rows):
    bg = ROW_ALT if i % 2 == 0 else WHITE
    for j, val in enumerate(row):
        style_data_cell(enroll_tbl.cell(i+1, j), val, font_size=7.5, fill_color=bg,
                       bold=(j==2), alignment=PP_ALIGN.LEFT if j < 2 else PP_ALIGN.CENTER)

enroll_col_widths = [Inches(0.9), Inches(0.85), Inches(0.75), Inches(0.7), Inches(0.7), Inches(0.7), Inches(0.8), Inches(1.0)]
for j, w in enumerate(enroll_col_widths):
    enroll_tbl.columns[j].width = w

# Total line
add_text_box(slide2, Inches(0.5), enroll_tbl_top + Inches(1.7), Inches(7.2), Inches(0.25),
             "TOTAL CURRENT: 5,263 students  |  Boys: 2,690 (51.1%)  |  Girls: 2,573 (48.9%)  |  281 Sections  |  Morning + Afternoon shifts",
             font_size=8, bold=True, color=DARK_GREEN)

# ── Right panel: Scaling methodology ──
method_x = Inches(8.1)
add_text_box(slide2, method_x, enroll_y, Inches(7.4), Inches(0.3),
             "SCALING METHODOLOGY  (Proportional Distribution from Current Baseline)",
             font_size=11, bold=True, color=DARK_GREEN)

method_lines = [
    "STEP 1:  Current distribution ratios derived from actual 5,263 enrollment",
    "           Early Years 15.2% | Primary 31.1% | Intermediate 39.0% | Secondary 14.7%",
    "",
    "STEP 2:  Apply ratios to target capacity \u2192 derive student count per segment",
    "",
    "STEP 3:  Classrooms = Students \u00f7 25 (TBC max/class) + operational buffer (~8%)",
    "",
    "STEP 4:  NET area = Classrooms \u00d7 TBC unit areas + specialist spaces (labs, LRC, etc.)",
    "           Specialist ratio: 1 Science Lab per 10 classrooms; 1 ICT Lab per ~15 classes",
    "",
    "STEP 5:  GROSS (BUA) = NET \u00d7 grossing factors (1.45 academic / 1.65 high-service)",
    "",
    "STEP 6:  Gender separation from G2 onwards doubles classroom wings (Boys + Girls)",
]
add_rect(slide2, method_x, enroll_y + Inches(0.35), Inches(7.4), Inches(2.8), RGBColor(0xF1, 0xF8, 0xE9))
add_multiline_box(slide2, method_x + Inches(0.15), enroll_y + Inches(0.45),
                  Inches(7.1), Inches(2.6), method_lines, font_size=8,
                  color=DARK_GREY, bold_first=False, line_spacing=1.25)

# ── MAIN COMPARISON TABLE ──
comp_y = Inches(4.55)
add_text_box(slide2, Inches(0.5), comp_y, Inches(15), Inches(0.3),
             "THREE-SCENARIO CAPACITY COMPARISON  |  NET-First Computation to BUA & Cost",
             font_size=12, bold=True, color=DARK_GREEN)

# Compute scenarios
# Ratios from current: EY=15.2%, Primary=31.1%, Intermediate=39.0%, Secondary=14.7%
scenarios = {
    5500: {"label": "Scenario A\n5,500 Students"},
    6000: {"label": "Scenario B\n6,000 Students"},
    7000: {"label": "Scenario C\n7,000 Students\n(BoD Target)"},
}

def compute_scenario(total):
    ey = round(total * 0.152)
    pri = round(total * 0.311)
    inter = round(total * 0.390)
    sec = total - ey - pri - inter

    # Classrooms (25/class + 8% buffer)
    ey_cls = math.ceil(ey / 22 * 1.08)  # EY smaller class sizes
    pri_cls = math.ceil(pri / 25 * 1.08)
    inter_cls = math.ceil(inter / 25 * 1.08)
    sec_cls = math.ceil(sec / 25 * 1.08)
    total_cls = ey_cls + pri_cls + inter_cls + sec_cls

    # NET areas (from space program)
    # Teaching: EY=62.5 avg, Primary=43.12, Inter/Sec=43.12
    teaching_net = ey_cls * 55.0 + pri_cls * 43.12 + inter_cls * 43.12 + sec_cls * 43.12

    # Labs: 1 per 10 classrooms, ~62 m² each
    n_labs = math.ceil((pri_cls + inter_cls + sec_cls) / 10) * 2  # boys+girls
    labs_net = n_labs * 65

    # ICT labs
    n_ict = math.ceil((pri_cls + inter_cls + sec_cls) / 15) * 2
    ict_net = n_ict * 65

    # Fixed shared facilities (scale slightly with student count)
    scale = total / 7000
    sen_net = round(736 * scale)
    admin_net = round(672 * max(scale, 0.85))  # admin doesn't shrink much
    staff_net = round(950 * scale)
    it_ops_net = round(641 * max(scale, 0.85))
    food_net = round(3380 * scale)
    sports_net = round(3981 * max(scale, 0.8))
    audit_commons_net = round(4070 * max(scale, 0.8))

    total_net = teaching_net + labs_net + ict_net + sen_net + admin_net + staff_net + it_ops_net + food_net + sports_net + audit_commons_net

    # Grossing
    academic_gross = (teaching_net + labs_net + ict_net + sen_net + admin_net + staff_net) * 1.45
    high_service_gross = (food_net + sports_net + audit_commons_net) * 1.65
    ops_gross = it_ops_net * 1.55
    total_gross = academic_gross + high_service_gross + ops_gross

    # Cost (SAR 4,580-4,960 per m² all-in mid-level)
    cost_low = round(total_gross * 4580 / 1e6)
    cost_high = round(total_gross * 4960 / 1e6)

    # Footprint
    footprint = round(total_gross / 3)
    coverage = round(footprint / 25000 * 100)

    return {
        "ey": ey, "pri": pri, "inter": inter, "sec": sec,
        "ey_cls": ey_cls, "pri_cls": pri_cls, "inter_cls": inter_cls, "sec_cls": sec_cls,
        "total_cls": total_cls,
        "teaching_net": round(teaching_net),
        "labs_net": labs_net + ict_net,
        "support_net": sen_net + admin_net + staff_net + it_ops_net,
        "shared_net": food_net + sports_net + audit_commons_net,
        "total_net": round(total_net),
        "total_gross": round(total_gross),
        "footprint": footprint,
        "coverage": coverage,
        "cost_low": cost_low,
        "cost_high": cost_high,
    }

s5500 = compute_scenario(5500)
s6000 = compute_scenario(6000)
s7000 = compute_scenario(7000)

comp_tbl_top = comp_y + Inches(0.35)
comp_tbl = add_table(slide2, 16, 5, Inches(0.5), comp_tbl_top, Inches(15), Inches(3.7))

# Headers
comp_headers = ["PARAMETER", "UNIT", "SCENARIO A\n5,500 Students", "SCENARIO B\n6,000 Students", "SCENARIO C\n7,000 Students\n(BoD Target)"]
for j, h in enumerate(comp_headers):
    style_header_cell(comp_tbl.cell(0, j), h, font_size=8)

def fmt_k(n):
    if n >= 1000:
        return f"{n:,}"
    return str(n)

comp_data = [
    # STUDENTS
    ["STUDENT DISTRIBUTION", "", "", "", ""],
    ["  Early Years (Nursery\u2013KG)", "students", str(s5500['ey']), str(s6000['ey']), str(s7000['ey'])],
    ["  Primary (G1\u2013G4)", "students", str(s5500['pri']), str(s6000['pri']), str(s7000['pri'])],
    ["  Intermediate (G5\u2013G9)", "students", str(s5500['inter']), str(s6000['inter']), str(s7000['inter'])],
    ["  Secondary (G10\u2013G12)", "students", str(s5500['sec']), str(s6000['sec']), str(s7000['sec'])],
    # CLASSROOMS
    ["TOTAL CLASSROOMS (incl. buffer)", "rooms", str(s5500['total_cls']), str(s6000['total_cls']), str(s7000['total_cls'])],
    # AREAS
    ["AREA BUILD-UP", "", "", "", ""],
    ["  Teaching Spaces NET", "m\u00b2", fmt_k(s5500['teaching_net']), fmt_k(s6000['teaching_net']), fmt_k(s7000['teaching_net'])],
    ["  Labs & Specialist NET", "m\u00b2", fmt_k(s5500['labs_net']), fmt_k(s6000['labs_net']), fmt_k(s7000['labs_net'])],
    ["  Support (SEN/Admin/Staff/IT)", "m\u00b2", fmt_k(s5500['support_net']), fmt_k(s6000['support_net']), fmt_k(s7000['support_net'])],
    ["  Shared (Food/Sports/Audit)", "m\u00b2", fmt_k(s5500['shared_net']), fmt_k(s6000['shared_net']), fmt_k(s7000['shared_net'])],
    ["  TOTAL NET AREA", "m\u00b2", fmt_k(s5500['total_net']), fmt_k(s6000['total_net']), fmt_k(s7000['total_net'])],
    ["  TOTAL GROSS / BUA", "m\u00b2", fmt_k(s5500['total_gross']), fmt_k(s6000['total_gross']), fmt_k(s7000['total_gross'])],
    # SITE
    ["SITE COVERAGE (3 floors)", "%", f"{s5500['coverage']}%", f"{s6000['coverage']}%", f"{s7000['coverage']}%"],
    # COST
    ["EST. CONSTRUCTION COST (mid)", "SAR M", f"{s5500['cost_low']}\u2013{s5500['cost_high']}M", f"{s6000['cost_low']}\u2013{s6000['cost_high']}M", f"{s7000['cost_low']}\u2013{s7000['cost_high']}M"],
]

section_rows = [0, 6]  # rows that are section headers
highlight_rows = [5, 11, 12, 14]  # key totals

for i, row in enumerate(comp_data):
    is_section = i in section_rows
    is_highlight = i in highlight_rows
    for j, val in enumerate(row):
        if is_section:
            style_data_cell(comp_tbl.cell(i+1, j), val, font_size=7.5, bold=True,
                           fill_color=RGBColor(0xE3, 0xF2, 0xFD), font_color=DARK_GREEN,
                           alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        elif is_highlight:
            style_data_cell(comp_tbl.cell(i+1, j), val, font_size=8, bold=True,
                           fill_color=RGBColor(0xFF, 0xF8, 0xE1), font_color=DARK_GREEN,
                           alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        else:
            bg = ROW_ALT if i % 2 == 0 else WHITE
            style_data_cell(comp_tbl.cell(i+1, j), val, font_size=7.5,
                           fill_color=bg,
                           alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

comp_col_widths = [Inches(3.2), Inches(0.8), Inches(3.2), Inches(3.2), Inches(4.6)]
for j, w in enumerate(comp_col_widths):
    comp_tbl.columns[j].width = w

# Footer
add_rect(slide2, Inches(0), Inches(8.55), Inches(16), Inches(0.45), DARK_GREEN)
add_text_box(slide2, Inches(0.5), Inches(8.58), Inches(10), Inches(0.35),
             "CONFIDENTIAL  |  Pakistan International School (English Section), Riyadh  |  Basis of Design v0.4",
             font_size=8, bold=False, color=GOLD)
add_text_box(slide2, Inches(12), Inches(8.58), Inches(3.5), Inches(0.35),
             "Prepared for Ambassador / SMC Briefing",
             font_size=8, bold=False, color=WHITE, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – FACILITY REQUIREMENTS & TIMELINE
# ═══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

# Top banner
add_rect(slide3, Inches(0), Inches(0), Inches(16), Inches(1.15), DARK_GREEN)
add_text_box(slide3, Inches(0.5), Inches(0.15), Inches(10), Inches(0.55),
             "FACILITY REQUIREMENTS BY LEVEL  |  TIMELINE & COST SCENARIOS",
             font_size=22, bold=True, color=WHITE)
add_text_box(slide3, Inches(0.5), Inches(0.65), Inches(10), Inches(0.4),
             "TBC-Mandated Facilities per 10 Classrooms + Construction Delivery Strategy  |  B + G + 2 Configuration",
             font_size=11, bold=False, color=GOLD)
add_text_box(slide3, Inches(12.5), Inches(0.25), Inches(3), Inches(0.5),
             "SLIDE 3 OF 3", font_size=10, bold=True, color=GOLD, alignment=PP_ALIGN.RIGHT)
add_text_box(slide3, Inches(12.5), Inches(0.55), Inches(3), Inches(0.4),
             "Facility Matrix & Implementation Roadmap", font_size=10, bold=False, color=WHITE, alignment=PP_ALIGN.RIGHT)

# ── LEFT: Facility requirements per 10 classrooms ──
fac_y = Inches(1.35)
add_text_box(slide3, Inches(0.5), fac_y, Inches(8), Inches(0.3),
             "TBC-MANDATED FACILITIES PER EDUCATION LEVEL  (Category A, per 10 Classrooms)",
             font_size=11, bold=True, color=DARK_GREEN)

fac_tbl_top = fac_y + Inches(0.35)
fac_tbl = add_table(slide3, 9, 7, Inches(0.5), fac_tbl_top, Inches(8.8), Inches(3.0))

fac_headers = ["Facility Type", "Nursery", "KG", "Elementary\n(Boys/Girls)", "Intermediate\n(Boys/Girls)", "Secondary\n(Boys/Girls)", "Basis"]
for j, h in enumerate(fac_headers):
    style_header_cell(fac_tbl.cell(0, j), h, font_size=7)

fac_data = [
    ["Classrooms (max 25/class)", "1.4\u20131.8 m\u00b2/ch", "2.5 m\u00b2/ch", "1.3 m\u00b2/st", "1.4 m\u00b2/st", "1.5 m\u00b2/st", "TBC"],
    ["Science Lab", "\u2014", "\u2014", "1 per 10 cls\n60.1 m\u00b2", "1 per 10 cls\n62.96 m\u00b2", "1 per 10 cls\n69.9 m\u00b2", "TBC"],
    ["Computer/Language Lab", "\u2014", "\u2014", "60.1 m\u00b2", "62.96 m\u00b2", "69.9 m\u00b2", "TBC"],
    ["Arts Atelier / Vocational", "\u2014", "\u2014", "41.9 m\u00b2", "51.0 m\u00b2", "56.7 m\u00b2", "TBC"],
    ["Learning Resource Center", "\u2014", "Multi-media\nLibrary", "75.1 m\u00b2", "74.7 m\u00b2", "88.6 m\u00b2", "TBC"],
    ["Multi-Purpose Room", "\u2014", "\u2014", "42.5 m\u00b2", "45.9 m\u00b2", "51.0 m\u00b2", "TBC"],
    ["Outdoor Pitch/Court", "Half class\narea", "Double class\narea", "Min 300 m\u00b2\npitch", "Min 400 m\u00b2\ncourt", "Min 400 m\u00b2\ncourt", "TBC"],
    ["Toilets", "1.8 m\u00b2/15st", "1.8 m\u00b2/class", "1.35 m\u00b2/20st", "1.35 m\u00b2/20st", "1.35 m\u00b2/20st", "TBC"],
]

for i, row in enumerate(fac_data):
    bg = ROW_ALT if i % 2 == 0 else WHITE
    for j, val in enumerate(row):
        style_data_cell(fac_tbl.cell(i+1, j), val, font_size=7, fill_color=bg,
                       alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

fac_col_widths = [Inches(1.6), Inches(1.0), Inches(1.0), Inches(1.3), Inches(1.3), Inches(1.3), Inches(0.6)]
for j, w in enumerate(fac_col_widths):
    fac_tbl.columns[j].width = w

# ── Classroom count by scenario ──
cls_y = fac_tbl_top + Inches(3.15)
add_text_box(slide3, Inches(0.5), cls_y, Inches(8.5), Inches(0.3),
             "CLASSROOM & LAB COUNT BY SCENARIO",
             font_size=10, bold=True, color=DARK_GREEN)

cls_tbl_top = cls_y + Inches(0.3)
cls_tbl = add_table(slide3, 7, 5, Inches(0.5), cls_tbl_top, Inches(8.8), Inches(2.45))

cls_headers = ["Segment", "Metric", "5,500 Students", "6,000 Students", "7,000 Students"]
for j, h in enumerate(cls_headers):
    style_header_cell(cls_tbl.cell(0, j), h, font_size=7)

# Recompute individual segment classrooms
cls_data = [
    ["Early Years", "Classrooms", str(s5500['ey_cls']), str(s6000['ey_cls']), str(s7000['ey_cls'])],
    ["Primary (G1-G4)", "Classrooms", str(s5500['pri_cls']), str(s6000['pri_cls']), str(s7000['pri_cls'])],
    ["Intermediate (G5-G9)", "Classrooms", str(s5500['inter_cls']), str(s6000['inter_cls']), str(s7000['inter_cls'])],
    ["Secondary (G10-G12)", "Classrooms", str(s5500['sec_cls']), str(s6000['sec_cls']), str(s7000['sec_cls'])],
    ["TOTAL CLASSROOMS", "Rooms", str(s5500['total_cls']), str(s6000['total_cls']), str(s7000['total_cls'])],
    ["Science + ICT Labs", "Rooms",
     str(math.ceil((s5500['pri_cls']+s5500['inter_cls']+s5500['sec_cls'])/10)*2 + math.ceil((s5500['pri_cls']+s5500['inter_cls']+s5500['sec_cls'])/15)*2),
     str(math.ceil((s6000['pri_cls']+s6000['inter_cls']+s6000['sec_cls'])/10)*2 + math.ceil((s6000['pri_cls']+s6000['inter_cls']+s6000['sec_cls'])/15)*2),
     str(math.ceil((s7000['pri_cls']+s7000['inter_cls']+s7000['sec_cls'])/10)*2 + math.ceil((s7000['pri_cls']+s7000['inter_cls']+s7000['sec_cls'])/15)*2)],
]

for i, row in enumerate(cls_data):
    is_total = (i == 4)
    bg = RGBColor(0xFF, 0xF8, 0xE1) if is_total else (ROW_ALT if i % 2 == 0 else WHITE)
    for j, val in enumerate(row):
        style_data_cell(cls_tbl.cell(i+1, j), val, font_size=7.5, bold=is_total,
                       fill_color=bg, font_color=DARK_GREEN if is_total else BLACK,
                       alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

cls_col_widths = [Inches(1.6), Inches(0.8), Inches(2.0), Inches(2.0), Inches(2.4)]
for j, w in enumerate(cls_col_widths):
    cls_tbl.columns[j].width = w

# ── RIGHT: Timeline & Cost Scenarios ──
right_x = Inches(9.7)
add_text_box(slide3, right_x, fac_y, Inches(6), Inches(0.3),
             "IMPLEMENTATION TIMELINE  |  DESIGN-BID-BUILD MODEL",
             font_size=11, bold=True, color=DARK_GREEN)

# Timeline table
tl_top = fac_y + Inches(0.35)
tl_tbl = add_table(slide3, 9, 3, right_x, tl_top, Inches(5.8), Inches(2.9))

tl_headers = ["Stage", "Duration", "Cumulative"]
for j, h in enumerate(tl_headers):
    style_header_cell(tl_tbl.cell(0, j), h, font_size=8)

tl_data = [
    ["1. Basis of Design (BoD)", "2\u20133 months", "Month 3"],
    ["2. Concept Design", "3\u20134 months", "Month 7"],
    ["3. Schematic Design", "4\u20136 months", "Month 13"],
    ["4. Detailed Design / IFC", "4\u20135 months", "Month 18"],
    ["5. Tender & Contractor Award", "2\u20133 months", "Month 21"],
    ["6. Construction Phase", "18\u201322 months", "Month 43"],
    ["7. Handover & Commissioning", "2\u20133 months", "Month 46"],
    ["TOTAL PROJECT DURATION", "30\u201346 months", "~3\u20134 years"],
]

for i, row in enumerate(tl_data):
    is_total = (i == 7)
    bg = RGBColor(0xFF, 0xF8, 0xE1) if is_total else (ROW_ALT if i % 2 == 0 else WHITE)
    for j, val in enumerate(row):
        style_data_cell(tl_tbl.cell(i+1, j), val, font_size=7.5, bold=is_total,
                       fill_color=bg, font_color=DARK_GREEN if is_total else BLACK,
                       alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

tl_col_widths = [Inches(2.8), Inches(1.4), Inches(1.6)]
for j, w in enumerate(tl_col_widths):
    tl_tbl.columns[j].width = w

# ── Cost scenario comparison ──
cost_y = tl_top + Inches(3.15)
add_text_box(slide3, right_x, cost_y, Inches(5.8), Inches(0.3),
             "COST SCENARIO COMPARISON  (Excl. Land & Professional Fees)",
             font_size=10, bold=True, color=DARK_GREEN)

cost_tbl_top = cost_y + Inches(0.3)
cost_tbl = add_table(slide3, 5, 4, right_x, cost_tbl_top, Inches(5.8), Inches(1.8))

cost_headers = ["Scenario", "Spec Level", "Est. Range (SAR)", "OPEX Profile"]
for j, h in enumerate(cost_headers):
    style_header_cell(cost_tbl.cell(0, j), h, font_size=7)

cost_data = [
    ["Code Minimum", "VRF, reduced finish", "205\u2013220M", "HIGH"],
    ["Mid-Institutional\n(Adopted Baseline)", "CHW HVAC, mid finish\nFull ICT, AV included", "240\u2013260M", "BALANCED"],
    ["Enhanced Campus", "Premium facade, BMS\nAdvanced acoustics", "270\u2013295M", "LOW"],
    ["+ Contingency (7\u201310%)", "\u2014", "+17\u201330M", "\u2014"],
]

for i, row in enumerate(cost_data):
    is_baseline = (i == 1)
    bg = RGBColor(0xFF, 0xF8, 0xE1) if is_baseline else (ROW_ALT if i % 2 == 0 else WHITE)
    for j, val in enumerate(row):
        style_data_cell(cost_tbl.cell(i+1, j), val, font_size=7.5,
                       bold=is_baseline,
                       fill_color=bg,
                       font_color=DARK_GREEN if is_baseline else BLACK,
                       alignment=PP_ALIGN.LEFT if j < 2 else PP_ALIGN.CENTER)

cost_col_widths = [Inches(1.3), Inches(1.8), Inches(1.3), Inches(1.1)]
for j, w in enumerate(cost_col_widths):
    cost_tbl.columns[j].width = w

# ── Building Configuration callout ──
config_y = cost_tbl_top + Inches(2.0)
add_rect(slide3, right_x, config_y, Inches(5.8), Inches(0.95), RGBColor(0xF1, 0xF8, 0xE9))
config_lines = [
    "BUILDING CONFIGURATION:  Basement + Ground + Floor 1 + Floor 2  (B + G + 2)",
    "Plot: 25,000 m\u00b2  |  Footprint: ~17,500 m\u00b2/floor  |  Site Coverage: ~70%  |  Urban high-density model",
    "Basement: Staff parking + MEP plant + Fire tanks + Storage  |  Above grade: Academic + Shared functions",
    "Delivery: Traditional Design-Bid-Build  |  Phasing Option: Phase 1 (5,000 cap) + Phase 2 (2,000 expansion)",
]
add_multiline_box(slide3, right_x + Inches(0.15), config_y + Inches(0.08),
                  Inches(5.5), Inches(0.85), config_lines, font_size=7.5,
                  color=DARK_GREY, bold_first=True, line_spacing=1.35)

# ── Bottom key assumptions ──
ka_y = Inches(7.6)
add_rect(slide3, Inches(0.5), ka_y, Inches(15), Inches(0.75), RGBColor(0xFB, 0xE9, 0xE7))
ka_lines = [
    "KEY ASSUMPTIONS & NOTES FOR AMBASSADOR REVIEW:",
    "\u25cf All TBC-mandated areas comply with Category A (Riyadh) standards  |  \u25cf Max 25 students/classroom (TBC)  |  \u25cf Gender separation from Grade 2 per Saudi regulation",
    "\u25cf Non-TBC items (SEN, Food, Sports, Auditorium) follow international best practice  |  \u25cf Costs are 2025 SAR planning estimates excl. land, fees & inflation  |  \u25cf 7,000 is design target; current enrollment = 5,263",
    "\u25cf NET \u2192 GROSS factors: Academic 1.45\u00d7 / High-service 1.65\u00d7 / Ops 1.55\u00d7  |  \u25cf Structural & MEP designed for full 7,000 capacity regardless of phasing",
]
add_multiline_box(slide3, Inches(0.65), ka_y + Inches(0.05),
                  Inches(14.7), Inches(0.7), ka_lines, font_size=7,
                  color=ACCENT_RED, bold_first=True, line_spacing=1.3)

# Footer
add_rect(slide3, Inches(0), Inches(8.55), Inches(16), Inches(0.45), DARK_GREEN)
add_text_box(slide3, Inches(0.5), Inches(8.58), Inches(10), Inches(0.35),
             "CONFIDENTIAL  |  Pakistan International School (English Section), Riyadh  |  Basis of Design v0.4",
             font_size=8, bold=False, color=GOLD)
add_text_box(slide3, Inches(12), Inches(8.58), Inches(3.5), Inches(0.35),
             "Prepared for Ambassador / SMC Briefing",
             font_size=8, bold=False, color=WHITE, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
output_path = "/home/user/PISES/PISES_Ambassador_Highlights_Deck.pptx"
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
print(f"Format: 16:9 widescreen (16\" x 9\")")
