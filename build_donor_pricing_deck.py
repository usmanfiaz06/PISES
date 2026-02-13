#!/usr/bin/env python3
"""
PISES New Campus – Donor Unit Pricing PowerPoint Deck Generator
Produces a 5-slide executive donor briefing matching the xlsx data.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ─────────────────────────────────────────────────────────────────────────────
# COLOUR PALETTE
# ─────────────────────────────────────────────────────────────────────────────
DARK_GREEN   = RGBColor(0x01, 0x41, 0x1C)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GOLD         = RGBColor(0xC4, 0x9A, 0x2A)
LIGHT_GREEN  = RGBColor(0x4C, 0xAF, 0x50)
LIGHT_GREY   = RGBColor(0xF5, 0xF5, 0xF5)
MID_GREY     = RGBColor(0xBD, 0xBD, 0xBD)
DARK_GREY    = RGBColor(0x42, 0x42, 0x42)
BLACK        = RGBColor(0x00, 0x00, 0x00)
HEADER_BG    = RGBColor(0x01, 0x41, 0x1C)
ROW_ALT      = RGBColor(0xE8, 0xF5, 0xE9)
ACCENT_RED   = RGBColor(0xC6, 0x28, 0x28)
MED_GREEN    = RGBColor(0x38, 0x8E, 0x3C)
ACCENT_GOLD  = RGBColor(0xFF, 0xF8, 0xE1)
LIGHT_BG     = RGBColor(0xF1, 0xF8, 0xE9)
BLUE_ACCENT  = RGBColor(0x15, 0x65, 0xC0)

# ─────────────────────────────────────────────────────────────────────────────
# CONSTANTS (from build_donor_pricing.py)
# ─────────────────────────────────────────────────────────────────────────────
TOTAL_COST_SAR = 250_000_000
TOTAL_BUA = 52_400
COST_PER_BUA_M2 = TOTAL_COST_SAR / TOTAL_BUA
SAR_TO_USD = 1 / 3.75
GF_ACADEMIC = 1.45
GF_HIGH_SERVICE = 1.65
GF_OPERATIONS = 1.55

def cost_per_unit(net_m2, gf):
    return round(net_m2 * gf * COST_PER_BUA_M2)

def usd(sar):
    return round(sar * SAR_TO_USD)

def fmt_sar(n):
    if n >= 1_000_000:
        return f"SAR {n/1_000_000:,.1f}M"
    return f"SAR {n:,.0f}"

def fmt_usd(n):
    if n >= 1_000_000:
        return f"USD {n/1_000_000:,.1f}M"
    return f"USD {n:,.0f}"

def fmt_both(sar):
    return f"{fmt_sar(sar)} / {fmt_usd(usd(sar))}"

# ─────────────────────────────────────────────────────────────────────────────
# PRESENTATION SETUP
# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)
TOTAL_SLIDES = 5

# ─────────────────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────────────────
def add_bg(slide, color=DARK_GREEN):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=11,
                      color=BLACK, font_name='Calibri', alignment=PP_ALIGN.LEFT,
                      bold_first=False, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.font.bold = (bold_first and i == 0)
        p.space_after = Pt(2)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_table(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, left, top, width, height).table

def style_header_cell(cell, text, font_size=9):
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = HEADER_BG
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def style_data_cell(cell, text, font_size=8, bold=False, alignment=PP_ALIGN.CENTER,
                    fill_color=None, font_color=BLACK):
    cell.text = str(text)
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color
        p.font.name = 'Calibri'
        p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_banner(slide, slide_num, title, subtitle, tag_line):
    add_rect(slide, Inches(0), Inches(0), Inches(16), Inches(1.15), DARK_GREEN)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.55),
                 title, font_size=22, bold=True, color=WHITE)
    add_text_box(slide, Inches(0.5), Inches(0.65), Inches(10), Inches(0.4),
                 subtitle, font_size=11, bold=False, color=GOLD)
    add_text_box(slide, Inches(12.5), Inches(0.25), Inches(3), Inches(0.5),
                 f"SLIDE {slide_num} OF {TOTAL_SLIDES}", font_size=10, bold=True,
                 color=GOLD, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(12.5), Inches(0.55), Inches(3), Inches(0.4),
                 tag_line, font_size=10, bold=False, color=WHITE, alignment=PP_ALIGN.RIGHT)

def add_footer(slide):
    add_rect(slide, Inches(0), Inches(8.55), Inches(16), Inches(0.45), DARK_GREEN)
    add_text_box(slide, Inches(0.5), Inches(8.58), Inches(10), Inches(0.35),
                 "CONFIDENTIAL  |  Pakistan International School (English Section), Riyadh  |  Donor Unit Pricing v1.0",
                 font_size=8, bold=False, color=GOLD)
    add_text_box(slide, Inches(12), Inches(8.58), Inches(3.5), Inches(0.35),
                 "Prepared for Donor / SMC Briefing",
                 font_size=8, bold=False, color=WHITE, alignment=PP_ALIGN.RIGHT)

def add_kpi_card(slide, left, top, width, height, label, value, sub="",
                 bg_color=WHITE, value_color=DARK_GREEN, label_color=DARK_GREY):
    add_rect(slide, left, top, width, height, bg_color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.12),
                 width - Inches(0.3), Inches(0.55),
                 value, font_size=22, bold=True, color=value_color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.62),
                 width - Inches(0.3), Inches(0.35),
                 label, font_size=9, bold=False, color=label_color, alignment=PP_ALIGN.CENTER)
    if sub:
        add_text_box(slide, left + Inches(0.15), top + Inches(0.88),
                     width - Inches(0.3), Inches(0.25),
                     sub, font_size=7, bold=False, color=MID_GREY, alignment=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# CATEGORY DATA (matches xlsx Category Summary sheet)
# ─────────────────────────────────────────────────────────────────────────────
CATEGORIES = [
    ("Classrooms & Teaching", 320, 14749, 102_035_420, 39.3),
    ("Science Laboratories", 40, 2193, 17_262_169, 6.6),
    ("Computer & ICT Labs", 16, 1060, 8_341_310, 3.2),
    ("Specialist Studios", 10, 748, 5_725_286, 2.2),
    ("Libraries & LRC", 6, 477, 3_298_472, 1.3),
    ("Sports & PE", 17, 6957, 54_766_460, 21.1),
    ("Dining & Food", 8, 2950, 23_151_240, 8.9),
    ("Auditorium & Assembly", 14, 3120, 24_198_469, 9.3),
    ("Exam Centre", 3, 870, 6_018_606, 2.3),
    ("SEN & Wellbeing", 34, 606, 4_192_264, 1.6),
    ("Staff & PD", 16, 810, 5_603_528, 2.2),
    ("Administration", 10, 356, 2_462_786, 0.9),
    ("IT & Security", 5, 96, 709_925, 0.3),
    ("Transport", 1, 43, 317_987, 0.1),
    ("Prayer Spaces", 4, 240, 1_660_304, 0.6),
]

GRAND_TOTAL_SAR = 259_744_226
GRAND_TOTAL_USD = 69_265_127


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE & PROJECT OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, DARK_GREEN)

# Central title block
add_text_box(slide1, Inches(1), Inches(1.5), Inches(14), Inches(1.0),
             "PISES NEW CAMPUS", font_size=44, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(1), Inches(2.5), Inches(14), Inches(0.7),
             "UNIT-BASED DONOR PRICING DECK", font_size=28, bold=True, color=GOLD,
             alignment=PP_ALIGN.CENTER)

# Divider line
add_rect(slide1, Inches(5), Inches(3.4), Inches(6), Inches(0.04), GOLD)

add_text_box(slide1, Inches(1), Inches(3.7), Inches(14), Inches(0.5),
             "Pakistan International School (English Section), Riyadh  |  Al Safa Plot  |  25,000 m\u00b2",
             font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)

# KPI row
kpi_y = Inches(4.6)
kpi_h = Inches(1.3)
kpi_w = Inches(2.6)
gap = Inches(0.35)
start_x = Inches(1.25)

kpis = [
    ("SAR 250M", "Total Project Cost", "Mid-Institutional Spec"),
    ("7,000", "Student Capacity", "Design Target"),
    ("52,400 m\u00b2", "Total Built-Up Area", "NET \u00d7 Grossing Factors"),
    ("504", "Total Donor Units", "58 Unique Unit Types"),
    ("SAR 260M", "Sum of All Units", "Unit-level detail pricing"),
]
for i, (val, label, sub) in enumerate(kpis):
    x = start_x + i * (kpi_w + gap)
    add_kpi_card(slide1, x, kpi_y, kpi_w, kpi_h, label, val, sub,
                 bg_color=RGBColor(0x1B, 0x5E, 0x20), value_color=GOLD,
                 label_color=WHITE)

# Bottom info
add_text_box(slide1, Inches(1), Inches(6.5), Inches(14), Inches(0.4),
             "1 USD = 3.75 SAR  |  Prices include construction, MEP, fit-out, ICT & furniture  |  Excluding land, professional fees & inflation",
             font_size=11, bold=False, color=MID_GREY, alignment=PP_ALIGN.CENTER)

add_text_box(slide1, Inches(1), Inches(7.1), Inches(14), Inches(0.4),
             "All facilities comply with Saudi Building Code 2024 and TBC Category A Standards",
             font_size=10, bold=False, color=MID_GREY, alignment=PP_ALIGN.CENTER)

# Footer
add_rect(slide1, Inches(0), Inches(8.55), Inches(16), Inches(0.45), RGBColor(0x00, 0x2E, 0x14))
add_text_box(slide1, Inches(0.5), Inches(8.58), Inches(10), Inches(0.35),
             "CONFIDENTIAL  |  PISES Donor Unit Pricing v1.0  |  2025",
             font_size=8, bold=False, color=GOLD)
add_text_box(slide1, Inches(12), Inches(8.58), Inches(3.5), Inches(0.35),
             f"SLIDE 1 OF {TOTAL_SLIDES}",
             font_size=8, bold=False, color=WHITE, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – COST SUMMARY BY CATEGORY
# ═══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)
add_banner(slide2, 2,
           "COST SUMMARY BY CATEGORY  |  15 FACILITY GROUPS",
           "High-level budget overview for donor briefings  |  7,000-Student Campus  |  SAR 250M Project",
           "Category Summary")

# Main table
tbl_top = Inches(1.45)
tbl = add_table(slide2, 17, 7, Inches(0.5), tbl_top, Inches(10.5), Inches(6.5))

headers = ["#", "Category", "Units", "Total NET m\u00b2", "Total Cost (SAR)", "Total Cost (USD)", "% of Budget"]
for j, h in enumerate(headers):
    style_header_cell(tbl.cell(0, j), h, font_size=8)

for i, (cat, units, net, cost, pct) in enumerate(CATEGORIES):
    bg = ROW_ALT if i % 2 == 0 else WHITE
    is_big = pct >= 5.0  # Highlight major categories
    vals = [str(i+1), cat, str(units), f"{net:,}", f"{cost:,}", f"{usd(cost):,}", f"{pct}%"]
    for j, val in enumerate(vals):
        fc = DARK_GREEN if is_big and j in (4, 5) else BLACK
        style_data_cell(tbl.cell(i+1, j), val, font_size=8,
                       fill_color=bg, bold=(is_big and j >= 4),
                       font_color=fc,
                       alignment=PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER)

# Grand total row
gt_row = 16
gt_vals = ["", "GRAND TOTAL", "504", "35,274", f"{GRAND_TOTAL_SAR:,}", f"{GRAND_TOTAL_USD:,}", "100%"]
for j, val in enumerate(gt_vals):
    style_data_cell(tbl.cell(gt_row, j), val, font_size=9, bold=True,
                   fill_color=DARK_GREEN, font_color=WHITE,
                   alignment=PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER)

col_widths = [Inches(0.4), Inches(2.4), Inches(0.7), Inches(1.3), Inches(2.0), Inches(2.0), Inches(1.2)]
for j, w in enumerate(col_widths):
    tbl.columns[j].width = w

# Right panel: Top 5 breakdown
panel_x = Inches(11.4)
add_text_box(slide2, panel_x, Inches(1.45), Inches(4.2), Inches(0.3),
             "TOP 5 CATEGORIES (82.6% of budget)",
             font_size=11, bold=True, color=DARK_GREEN)

top5 = [
    ("Classrooms", "SAR 102.0M", "39.3%", "320 rooms"),
    ("Sports & PE", "SAR 54.8M", "21.1%", "17 facilities"),
    ("Auditorium", "SAR 24.2M", "9.3%", "14 spaces"),
    ("Dining & Food", "SAR 23.2M", "8.9%", "8 facilities"),
    ("Science Labs", "SAR 17.3M", "6.6%", "40 labs"),
]

for i, (name, cost_str, pct_str, qty_str) in enumerate(top5):
    y = Inches(1.95) + i * Inches(0.95)
    add_rect(slide2, panel_x, y, Inches(4.2), Inches(0.82), LIGHT_BG)
    add_text_box(slide2, panel_x + Inches(0.15), y + Inches(0.05),
                 Inches(2.5), Inches(0.3),
                 name, font_size=11, bold=True, color=DARK_GREEN)
    add_text_box(slide2, panel_x + Inches(0.15), y + Inches(0.33),
                 Inches(2.0), Inches(0.25),
                 cost_str, font_size=10, bold=True, color=BLACK)
    add_text_box(slide2, panel_x + Inches(2.3), y + Inches(0.33),
                 Inches(1.8), Inches(0.25),
                 f"{pct_str}  |  {qty_str}", font_size=8, bold=False, color=DARK_GREY)
    # Percentage bar
    bar_w = float(pct_str.replace('%', '')) / 40.0 * 3.9
    add_rect(slide2, panel_x + Inches(0.15), y + Inches(0.62),
             Inches(bar_w), Inches(0.1), MED_GREEN)
    add_rect(slide2, panel_x + Inches(0.15) + Inches(bar_w), y + Inches(0.62),
             Inches(3.9 - bar_w), Inches(0.1), RGBColor(0xE0, 0xE0, 0xE0))

# Notes
add_rect(slide2, panel_x, Inches(6.85), Inches(4.2), Inches(1.15), ACCENT_GOLD)
note_lines = [
    "KEY NOTES:",
    "\u2022 Grand total (SAR 260M) reflects sum of all",
    "  individual units at planning-level estimates",
    "\u2022 Full campus mid-range: SAR 240\u2013260M",
    "\u2022 Grossing: Academic 1.45\u00d7 / Service 1.65\u00d7",
    "\u2022 All prices in 2025 SAR baseline",
]
add_multiline_box(slide2, panel_x + Inches(0.12), Inches(6.9),
                  Inches(4.0), Inches(1.05), note_lines, font_size=7.5,
                  color=DARK_GREY, bold_first=True, line_spacing=1.25)

add_footer(slide2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – TOP UNIT PRICING (Key Items from Unit Pricing sheet)
# ═══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)
add_banner(slide3, 3,
           "UNIT PRICING  |  KEY FACILITIES WITH COST PER UNIT",
           "58 unique unit types  |  NET m\u00b2 \u00d7 Grossing Factor \u00d7 SAR 4,771/m\u00b2 BUA  |  Full list in Excel workbook",
           "Unit-Level Detail")

# Selected high-interest units for the slide (from xlsx data)
units_data = [
    # (name, qty, net_m2, cost_unit_sar, total_sar, students)
    ("Standard Classroom (G1\u201312)", 249, 43.1, 298_302, 74_277_198, "25 students/room"),
    ("KG Classroom", 26, 62.5, 432_371, 11_241_646, "25 children/room"),
    ("Nursery Activity Room", 9, 45.0, 311_307, 2_801_763, "20\u201325 children"),
    ("Reception Classroom", 20, 62.5, 432_371, 8_647_420, "25 children/room"),
    ("Primary Science Lab", 13, 60.1, 473_115, 6_150_495, "25 students/lab"),
    ("Intermediate Science Lab", 7, 63.0, 495_630, 3_469_410, "25 students/lab"),
    ("Secondary Science Lab", 12, 69.9, 550_262, 6_603_144, "25 students/lab"),
    ("Primary Computer Lab", 6, 60.1, 473_115, 2_838_690, "25 students/session"),
    ("Secondary Computer Lab", 10, 69.9, 550_262, 5_502_620, "25 students/session"),
    ("Maker / Robotics Lab", 2, 120.0, 944_656, 1_889_312, "25 students/session"),
    ("Art Studio", 2, 90.0, 708_492, 1_416_984, "25 students/session"),
    ("Music / Drama Room", 2, 80.0, 629_771, 1_259_542, "30\u201340 students"),
    ("Primary Library / LRC", 2, 75.1, 519_537, 1_039_074, "40\u201360 students"),
    ("Secondary LRC", 2, 88.6, 612_929, 1_225_858, "50\u201370 students"),
    ("Sports Hall", 2, 900.0, 7_084_924, 14_169_848, "200+ students/day"),
    ("25m Swimming Pool", 1, 1717.0, 13_516_460, 13_516_460, "300+ students/wk"),
    ("Dining Hall (700-seat)", 2, 1100.0, 8_659_351, 17_318_702, "700 students/sitting"),
    ("Auditorium (300 seats)", 1, 740.0, 5_825_382, 5_825_382, "300-seat events"),
    ("Atrium / Learning Commons", 1, 2000.0, 15_744_275, 15_744_275, "2,000+ students"),
    ("Exam Hall (300 cands)", 1, 750.0, 5_188_454, 5_188_454, "300 candidates"),
    ("SEN Resource Room", 10, 25.0, 172_948, 1_729_480, "4\u20138 students"),
    ("Prayer Room / Musalla", 4, 60.0, 415_076, 1_660_304, "100\u2013150 per room"),
]

# Split into two columns
half = len(units_data) // 2
left_data = units_data[:half]
right_data = units_data[half:]

def draw_unit_table(slide, data, left_x, top_y, table_width):
    n_rows = len(data) + 1
    tbl = add_table(slide, n_rows, 6, left_x, top_y, table_width, Inches(0.3 * n_rows))
    hdrs = ["Unit Name", "Qty", "NET m\u00b2", "Cost/Unit (SAR)", "Total (SAR)", "Impact"]
    for j, h in enumerate(hdrs):
        style_header_cell(tbl.cell(0, j), h, font_size=7)

    for i, (name, qty, net, cpu, total, impact) in enumerate(data):
        bg = ROW_ALT if i % 2 == 0 else WHITE
        vals = [name, str(qty), f"{net:.0f}", f"{cpu:,}", f"{total:,}", impact]
        for j, val in enumerate(vals):
            al = PP_ALIGN.LEFT if j in (0, 5) else PP_ALIGN.CENTER
            if j in (3, 4):
                al = PP_ALIGN.RIGHT
            style_data_cell(tbl.cell(i+1, j), val, font_size=7, fill_color=bg,
                           alignment=al)

    # Column widths
    cw = [Inches(table_width / Inches(1) * r) for r in [0.28, 0.06, 0.09, 0.19, 0.19, 0.19]]
    for j, w in enumerate(cw):
        tbl.columns[j].width = int(w)
    return tbl

tbl_top = Inches(1.45)
draw_unit_table(slide3, left_data, Inches(0.3), tbl_top, Inches(7.6))
draw_unit_table(slide3, right_data, Inches(8.2), tbl_top, Inches(7.6))

# Bottom note
add_rect(slide3, Inches(0.3), Inches(8.0), Inches(15.4), Inches(0.35), ACCENT_GOLD)
add_text_box(slide3, Inches(0.5), Inches(8.03), Inches(15), Inches(0.3),
             "Full pricing for all 58 unit types (504 total units) available in the Excel workbook  |  "
             "Cost = NET m\u00b2 \u00d7 Grossing Factor \u00d7 SAR 4,771/m\u00b2 BUA  |  "
             "Grossing: Academic 1.45\u00d7 / High-Service 1.65\u00d7 / Operations 1.55\u00d7",
             font_size=8, bold=False, color=DARK_GREY, alignment=PP_ALIGN.CENTER)

add_footer(slide3)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – DONOR PACKAGES
# ═══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)
add_banner(slide4, 4,
           "DONOR PACKAGES  |  THREE GIVING TIERS",
           "Suggested giving levels with naming recognition  |  All amounts in SAR & USD",
           "Giving Opportunities")

# TIER 1: Individual Impact Gifts
tier1_y = Inches(1.45)
add_rect(slide4, Inches(0.5), tier1_y, Inches(4.7), Inches(0.4), MED_GREEN)
add_text_box(slide4, Inches(0.6), tier1_y + Inches(0.05), Inches(4.5), Inches(0.3),
             "INDIVIDUAL IMPACT GIFTS  (SAR 50K \u2013 500K)",
             font_size=11, bold=True, color=WHITE)

tier1_items = [
    ("Name a Classroom", "SAR 298,302", "USD 79,547", "25 students"),
    ("Equip a Science Lab", "SAR 473,115", "USD 126,164", "25 students/lab"),
    ("Build a Sensory Room", "SAR 166,031", "USD 44,275", "SEN students"),
    ("Sponsor a Library Corner", "SAR 519,537", "USD 138,543", "40\u201360 students"),
    ("Create an Art Atelier", "SAR 289,862", "USD 77,297", "25 young artists"),
]

t1_tbl = add_table(slide4, 6, 5, Inches(0.5), tier1_y + Inches(0.5), Inches(4.7), Inches(2.1))
for j, h in enumerate(["#", "Package", "SAR", "USD", "Impact"]):
    style_header_cell(t1_tbl.cell(0, j), h, font_size=7)
for i, (pkg, sar, usd_val, impact) in enumerate(tier1_items):
    bg = ROW_ALT if i % 2 == 0 else WHITE
    for j, val in enumerate([str(i+1), pkg, sar, usd_val, impact]):
        al = PP_ALIGN.LEFT if j in (1, 4) else PP_ALIGN.CENTER
        if j in (2, 3): al = PP_ALIGN.RIGHT
        style_data_cell(t1_tbl.cell(i+1, j), val, font_size=7.5, fill_color=bg, alignment=al)
t1_cw = [Inches(0.3), Inches(1.6), Inches(1.0), Inches(0.9), Inches(0.9)]
for j, w in enumerate(t1_cw):
    t1_tbl.columns[j].width = w

# TIER 2: Major Gifts
tier2_y = Inches(1.45)
tier2_x = Inches(5.6)
add_rect(slide4, tier2_x, tier2_y, Inches(4.7), Inches(0.4), DARK_GREEN)
add_text_box(slide4, tier2_x + Inches(0.1), tier2_y + Inches(0.05), Inches(4.5), Inches(0.3),
             "MAJOR GIFTS  (SAR 500K \u2013 5M)",
             font_size=11, bold=True, color=WHITE)

tier2_items = [
    ("Robotics Innovation Hub", "SAR 944,656", "USD 251,908", "STEM education"),
    ("Auditorium Naming", "SAR 5.8M", "USD 1.6M", "300-seat events"),
    ("Sports Hall Sponsor", "SAR 7.1M", "USD 1.9M", "200+ students/day"),
    ("Dining Experience", "SAR 8.7M", "USD 2.3M", "700 students/sitting"),
    ("Classroom Block (10)", "SAR 3.0M", "USD 795K", "250 students"),
]

t2_tbl = add_table(slide4, 6, 5, tier2_x, tier2_y + Inches(0.5), Inches(4.7), Inches(2.1))
for j, h in enumerate(["#", "Package", "SAR", "USD", "Impact"]):
    style_header_cell(t2_tbl.cell(0, j), h, font_size=7)
for i, (pkg, sar, usd_val, impact) in enumerate(tier2_items):
    bg = ROW_ALT if i % 2 == 0 else WHITE
    for j, val in enumerate([str(i+1), pkg, sar, usd_val, impact]):
        al = PP_ALIGN.LEFT if j in (1, 4) else PP_ALIGN.CENTER
        if j in (2, 3): al = PP_ALIGN.RIGHT
        style_data_cell(t2_tbl.cell(i+1, j), val, font_size=7.5, fill_color=bg, alignment=al)
t2_cw = [Inches(0.3), Inches(1.6), Inches(1.0), Inches(0.9), Inches(0.9)]
for j, w in enumerate(t2_cw):
    t2_tbl.columns[j].width = w

# TIER 3: Landmark Gifts
tier3_y = Inches(1.45)
tier3_x = Inches(10.7)
add_rect(slide4, tier3_x, tier3_y, Inches(4.8), Inches(0.4), RGBColor(0xB7, 0x14, 0x1C))
add_text_box(slide4, tier3_x + Inches(0.1), tier3_y + Inches(0.05), Inches(4.6), Inches(0.3),
             "LANDMARK GIFTS  (SAR 5M+)",
             font_size=11, bold=True, color=WHITE)

tier3_items = [
    ("Swimming Pool Complex", "SAR 13.5M", "USD 3.6M", "300+ students/wk"),
    ("Exam Centre", "SAR 6.4M", "USD 1.7M", "300 candidates"),
    ("Learning Commons", "SAR 15.7M", "USD 4.2M", "2,000+ students"),
    ("Early Years Wing (55 rooms)", "SAR 26.6M", "USD 7.1M", "800+ children"),
    ("Complete SEN Suite (34 rooms)", "SAR 4.2M", "USD 1.1M", "500+ SEN students"),
]

t3_tbl = add_table(slide4, 6, 5, tier3_x, tier3_y + Inches(0.5), Inches(4.8), Inches(2.1))
for j, h in enumerate(["#", "Package", "SAR", "USD", "Impact"]):
    style_header_cell(t3_tbl.cell(0, j), h, font_size=7)
for i, (pkg, sar, usd_val, impact) in enumerate(tier3_items):
    bg = ROW_ALT if i % 2 == 0 else WHITE
    for j, val in enumerate([str(i+1), pkg, sar, usd_val, impact]):
        al = PP_ALIGN.LEFT if j in (1, 4) else PP_ALIGN.CENTER
        if j in (2, 3): al = PP_ALIGN.RIGHT
        style_data_cell(t3_tbl.cell(i+1, j), val, font_size=7.5, fill_color=bg, alignment=al)
t3_cw = [Inches(0.3), Inches(1.7), Inches(1.0), Inches(0.9), Inches(0.9)]
for j, w in enumerate(t3_cw):
    t3_tbl.columns[j].width = w

# How to Give section
how_y = Inches(4.3)
add_rect(slide4, Inches(0.5), how_y, Inches(15), Inches(0.4), DARK_GREEN)
add_text_box(slide4, Inches(0.6), how_y + Inches(0.05), Inches(14.8), Inches(0.3),
             "HOW TO GIVE  |  RECOGNITION & CO-SPONSORSHIP",
             font_size=12, bold=True, color=WHITE)

how_lines = [
    "\u2022  Donors may sponsor any unit individually or combine units for larger impact",
    "\u2022  Naming rights available for gifts of SAR 250,000 and above (recognition plaque on facility)",
    "\u2022  Co-sponsorship welcomed \u2014 multiple donors can share the cost of larger facilities",
    "\u2022  Contact the PISES Development Office for customized giving plans and recognition",
]

for i, line in enumerate(how_lines):
    y = how_y + Inches(0.5) + i * Inches(0.35)
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    add_rect(slide4, Inches(0.5), y, Inches(15), Inches(0.32), bg)
    add_text_box(slide4, Inches(0.7), y + Inches(0.03), Inches(14.6), Inches(0.26),
                 line, font_size=10, bold=False, color=DARK_GREY)

# Impact statement
impact_y = Inches(6.1)
add_rect(slide4, Inches(1.5), impact_y, Inches(13), Inches(1.2), ACCENT_GOLD)
add_text_box(slide4, Inches(2), impact_y + Inches(0.15), Inches(12), Inches(0.45),
             "\"EVERY CONTRIBUTION BUILDS A FUTURE\"",
             font_size=20, bold=True, color=DARK_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide4, Inches(2), impact_y + Inches(0.6), Inches(12), Inches(0.45),
             "From a single classroom (SAR 298K) to an entire early years wing (SAR 26.6M) \u2014 "
             "every donor gift directly builds the infrastructure that will educate 7,000 Pakistani children in Riyadh.",
             font_size=11, bold=False, color=DARK_GREY, alignment=PP_ALIGN.CENTER)

# Visual bar: what different amounts build
bar_y = Inches(7.6)
add_text_box(slide4, Inches(0.5), bar_y - Inches(0.3), Inches(5), Inches(0.25),
             "WHAT YOUR GIFT CAN BUILD:", font_size=9, bold=True, color=DARK_GREEN)

gift_levels = [
    ("SAR 83K", "Assessment\nRoom", 0.8),
    ("SAR 173K", "Breakout\nRoom", 1.1),
    ("SAR 298K", "Classroom", 1.5),
    ("SAR 473K", "Science\nLab", 2.0),
    ("SAR 945K", "Robotics\nLab", 2.8),
    ("SAR 3.0M", "10-Room\nBlock", 4.0),
    ("SAR 7.1M", "Sports\nHall", 5.5),
    ("SAR 13.5M", "Swimming\nPool", 7.0),
    ("SAR 15.7M", "Learning\nCommons", 7.5),
]

x_pos = Inches(0.5)
for amount, label, bar_h_factor in gift_levels:
    bar_h = Inches(bar_h_factor * 0.08)
    bar_bottom = bar_y + Inches(0.6)
    add_rect(slide4, x_pos, bar_bottom - bar_h, Inches(1.4), bar_h, MED_GREEN)
    add_text_box(slide4, x_pos, bar_bottom - bar_h - Inches(0.25), Inches(1.4), Inches(0.2),
                 amount, font_size=6, bold=True, color=DARK_GREEN, alignment=PP_ALIGN.CENTER)
    add_text_box(slide4, x_pos, bar_bottom + Inches(0.02), Inches(1.4), Inches(0.3),
                 label, font_size=6, bold=False, color=DARK_GREY, alignment=PP_ALIGN.CENTER)
    x_pos += Inches(1.6)

add_footer(slide4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – QUICK REFERENCE: WHAT YOUR GIFT CAN BUILD
# ═══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)
add_banner(slide5, 5,
           "WHAT YOUR GIFT CAN BUILD  |  QUICK REFERENCE",
           "At-a-glance pricing by giving level  |  All amounts include construction, fit-out, ICT & furniture",
           "Quick Reference Card")

# Giving bands (from Quick Reference sheet in xlsx)
bands = [
    ("SAR 50,000 \u2013 100,000", MED_GREEN, [
        ("SEN Assessment Room", "SAR 83,000 / USD 22,000"),
        ("Counsellor Room", "SAR 83,000 / USD 22,000"),
        ("Breakout Room", "SAR 173,000 / USD 46,000"),
        ("Medical Clinic", "SAR 138,000 / USD 37,000"),
    ]),
    ("SAR 100,000 \u2013 300,000", MED_GREEN, [
        ("Nursery Bedroom", "SAR 156,000 / USD 42,000"),
        ("SEN Resource Room", "SAR 173,000 / USD 46,000"),
        ("Primary Art Atelier", "SAR 290,000 / USD 77,000"),
        ("Standard Classroom", "SAR 298,000 / USD 79,000"),
    ]),
    ("SAR 300,000 \u2013 500,000", DARK_GREEN, [
        ("Nursery Activity Room", "SAR 311,000 / USD 83,000"),
        ("KG / Reception Classroom", "SAR 432,000 / USD 115,000"),
        ("Primary Science Lab", "SAR 473,000 / USD 126,000"),
        ("Primary Computer Lab", "SAR 473,000 / USD 126,000"),
    ]),
    ("SAR 500,000 \u2013 1,000,000", DARK_GREEN, [
        ("Secondary Science Lab", "SAR 551,000 / USD 147,000"),
        ("Secondary Computer Lab", "SAR 551,000 / USD 147,000"),
        ("Music / Drama Room", "SAR 630,000 / USD 168,000"),
        ("Art Studio", "SAR 709,000 / USD 189,000"),
        ("Early Years Learning Commons", "SAR 830,000 / USD 221,000"),
        ("Maker / Robotics Lab", "SAR 945,000 / USD 252,000"),
    ]),
    ("SAR 1,000,000 \u2013 5,000,000", RGBColor(0xB7, 0x14, 0x1C), [
        ("Prayer Room / Musalla (\u00d74)", "SAR 1,660,000 / USD 443,000"),
        ("Classroom Block (10 rooms)", "SAR 2,980,000 / USD 795,000"),
        ("Exam Hall (300 candidates)", "SAR 5,188,000 / USD 1,383,000"),
    ]),
    ("SAR 5,000,000+", RGBColor(0xB7, 0x14, 0x1C), [
        ("Indoor Sports Hall", "SAR 7,085,000 / USD 1,889,000"),
        ("Dining Hall + Kitchen", "SAR 11,020,000 / USD 2,939,000"),
        ("Swimming Pool Complex", "SAR 13,519,000 / USD 3,605,000"),
        ("Auditorium (300 seats)", "SAR 5,824,000 / USD 1,553,000"),
        ("Atrium / Learning Commons", "SAR 15,741,000 / USD 4,198,000"),
    ]),
]

# Lay out in 2 columns, 3 bands each
col_x = [Inches(0.5), Inches(8.2)]
col_bands = [bands[:3], bands[3:]]

for col_idx, col_data in enumerate(col_bands):
    x = col_x[col_idx]
    y = Inches(1.45)

    for band_name, band_color, items in col_data:
        # Band header
        add_rect(slide5, x, y, Inches(7.2), Inches(0.38), band_color)
        add_text_box(slide5, x + Inches(0.15), y + Inches(0.04), Inches(6.9), Inches(0.3),
                     band_name, font_size=11, bold=True, color=WHITE)
        y += Inches(0.42)

        for i, (item_name, item_price) in enumerate(items):
            bg = LIGHT_BG if i % 2 == 0 else WHITE
            add_rect(slide5, x, y, Inches(7.2), Inches(0.33), bg)
            add_text_box(slide5, x + Inches(0.2), y + Inches(0.04), Inches(3.0), Inches(0.25),
                         item_name, font_size=10, bold=False, color=BLACK)
            add_text_box(slide5, x + Inches(3.4), y + Inches(0.04), Inches(3.6), Inches(0.25),
                         item_price, font_size=10, bold=True, color=DARK_GREEN,
                         alignment=PP_ALIGN.RIGHT)
            y += Inches(0.35)

        y += Inches(0.2)

# Bottom callout
add_rect(slide5, Inches(1.5), Inches(7.6), Inches(13), Inches(0.75), ACCENT_GOLD)
callout_lines = [
    "EVERY CONTRIBUTION BUILDS A FUTURE  |  PISES NEW CAMPUS  |  7,000 STUDENTS",
    "Contact the PISES Development Office for customized giving plans  |  Naming rights for gifts SAR 250,000+  |  Co-sponsorship welcomed",
]
add_multiline_box(slide5, Inches(2), Inches(7.65), Inches(12), Inches(0.7),
                  callout_lines, font_size=10, color=DARK_GREEN,
                  bold_first=True, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

add_footer(slide5)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
output_path = "/home/user/PISES/PISES_Donor_Unit_Pricing_Deck.pptx"
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
print(f"Format: 16:9 widescreen (16\" x 9\")")
